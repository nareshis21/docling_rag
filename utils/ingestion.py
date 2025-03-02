import json
import time
import os
from pathlib import Path
from typing import Dict, Any, List
import chromadb

from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import (
    AcceleratorDevice,
    AcceleratorOptions,
    PdfPipelineOptions,
    TableFormerMode
)
from docling.document_converter import DocumentConverter, PdfFormatOption
from docling_core.transforms.chunker.hybrid_chunker import HybridChunker
from langchain_community.embeddings.fastembed import FastEmbedEmbeddings

from docx import Document  # DOCX support
from pptx import Presentation  # PPTX support
from bs4 import BeautifulSoup  # HTML support


class DocumentProcessor:
    def __init__(self):
        """Initialize document processor with necessary components"""
        self.setup_document_converter()
        self.embed_model = FastEmbedEmbeddings()
        self.client = chromadb.PersistentClient(path="chroma_db")  # Persistent Storage

    def setup_document_converter(self):
        """Configure document converter with advanced processing capabilities"""
        pipeline_options = PdfPipelineOptions()
        pipeline_options.do_ocr = True
        pipeline_options.do_table_structure = True
        pipeline_options.table_structure_options.do_cell_matching = True
        pipeline_options.ocr_options.lang = ["en"]
        pipeline_options.table_structure_options.mode = TableFormerMode.ACCURATE

        try:
            pipeline_options.accelerator_options = AcceleratorOptions(
                num_threads=8, device=AcceleratorDevice.MPS
            )
        except Exception:
            print("⚠️ MPS is not available. Falling back to CPU.")
            pipeline_options.accelerator_options = AcceleratorOptions(
                num_threads=8, device=AcceleratorDevice.CPU
            )

        self.converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(
                    pipeline_options=pipeline_options,
                    backend=PyPdfiumDocumentBackend
                )
            }
        )

    def extract_chunk_metadata(self, chunk) -> Dict[str, Any]:
        """Extract essential metadata from a chunk"""
        metadata = {
            "text": chunk.text.strip(),
            "headings": [],
            "page_info": None,
            "content_type": None
        }

        if hasattr(chunk, 'meta'):
            if hasattr(chunk.meta, 'headings') and chunk.meta.headings:
                metadata["headings"] = chunk.meta.headings

            if hasattr(chunk.meta, 'doc_items'):
                for item in chunk.meta.doc_items:
                    if hasattr(item, 'label'):
                        metadata["content_type"] = str(item.label)

                    if hasattr(item, 'prov') and item.prov:
                        for prov in item.prov:
                            if hasattr(prov, 'page_no'):
                                metadata["page_info"] = prov.page_no

        return metadata

    def extract_text_from_docx(self, docx_path: str) -> List[str]:
        """Extract text from a DOCX file"""
        doc = Document(docx_path)
        return [para.text.strip() for para in doc.paragraphs if para.text.strip()]

    def extract_text_from_pptx(self, pptx_path: str) -> List[str]:
        """Extract text from a PPTX file"""
        ppt = Presentation(pptx_path)
        slides_text = []
        for slide in ppt.slides:
            text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            if text.strip():
                slides_text.append(text.strip())
        return slides_text

    def extract_text_from_html(self, html_path: str) -> List[str]:
        """Extract text from an HTML file"""
        with open(html_path, "r", encoding="utf-8") as file:
            soup = BeautifulSoup(file, "html.parser")
        return [text.strip() for text in soup.stripped_strings if text.strip()]

    def extract_text_from_txt(self, txt_path: str) -> List[str]:
        """Extract text from a TXT file"""
        with open(txt_path, "r", encoding="utf-8") as file:
            lines = file.readlines()
        return [line.strip() for line in lines if line.strip()]

    def process_document(self, file_path: str):
        """Process document and create searchable index with metadata"""
        print(f"📄 Processing document: {file_path}")
        start_time = time.time()
        file_ext = Path(file_path).suffix.lower()

        if file_ext == ".pdf":
            result = self.converter.convert(file_path)
            doc = result.document
            chunker = HybridChunker(tokenizer="jinaai/jina-embeddings-v3")
            chunks = list(chunker.chunk(doc))

            processed_chunks = []
            for chunk in chunks:
                metadata = self.extract_chunk_metadata(chunk)
                processed_chunks.append(metadata)

        elif file_ext == ".docx":
            texts = self.extract_text_from_docx(file_path)
            processed_chunks = [{"text": text, "headings": [], "content_type": "DOCX"} for text in texts]

        elif file_ext == ".pptx":
            texts = self.extract_text_from_pptx(file_path)
            processed_chunks = [{"text": text, "headings": [], "content_type": "PPTX"} for text in texts]

        elif file_ext == ".html":
            texts = self.extract_text_from_html(file_path)
            processed_chunks = [{"text": text, "headings": [], "content_type": "HTML"} for text in texts]

        elif file_ext == ".txt":
            texts = self.extract_text_from_txt(file_path)
            processed_chunks = [{"text": text, "headings": [], "content_type": "TXT"} for text in texts]

        else:
            print(f"❌ Unsupported file format: {file_ext}")
            return None

        print("✅ Chunking completed. Creating vector database...")
        collection = self.client.get_or_create_collection(name="document_chunks")

        documents = []
        embeddings = []
        metadata_list = []
        ids = []

        for idx, chunk in enumerate(processed_chunks):
            text = chunk.get('text', '').strip()
            if not text:
                print(f"⚠️ Skipping empty chunk at index {idx}")
                continue  # Skip empty chunks

            embedding = self.embed_model.embed_documents([text])[0]  # ✅ Corrected method
            documents.append(text)
            embeddings.append(embedding)
            metadata_list.append({
                "headings": json.dumps(chunk.get('headings', [])),
                "content_type": chunk.get('content_type', None)
            })
            ids.append(str(idx))

        if documents:
            collection.add(
                ids=ids,
                embeddings=embeddings,
                documents=documents,
                metadatas=metadata_list
            )
            print(f"✅ Successfully added {len(documents)} chunks to the database.")

        processing_time = time.time() - start_time
        print(f"✅ Document processing completed in {processing_time:.2f} seconds")
        return collection
